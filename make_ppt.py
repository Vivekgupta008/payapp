"""
Generates PaytmOfflinePay_Hackathon.pptx
Run: python3 make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ─────────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x29, 0x70)
BLUE   = RGBColor(0x00, 0xB9, 0xF1)
LIGHT  = RGBColor(0xE8, 0xF4, 0xFD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
GREEN  = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
GREY   = RGBColor(0x60, 0x60, 0x60)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def add_multiline(slide, lines, x, y, w, h,
                  size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                  spacing=1.15, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return txb

def pill(slide, text, x, y, w=Inches(2.4), h=Inches(0.38),
         bg=NAVY, fg=WHITE, size=11):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, text, x, y, w, h, size=size, bold=True,
             color=fg, align=PP_ALIGN.CENTER)

def section_header(slide, title, subtitle=None):
    """Left accent bar + title"""
    add_rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=NAVY)
    add_text(slide, title, Inches(0.3), Inches(0.25), Inches(12.7), Inches(0.7),
             size=36, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.3), Inches(0.88), Inches(12.7), Inches(0.5),
                 size=15, color=GREY, align=PP_ALIGN.LEFT)
    # divider
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.03), fill=LIGHT)

def card(slide, x, y, w, h, title=None, body_lines=None,
         title_color=NAVY, body_color=DARK, icon=None, bg=WHITE,
         title_size=15, body_size=12.5):
    add_rect(slide, x, y, w, h, fill=bg, line=LIGHT, line_w=Pt(0.5))
    cy = y + Inches(0.18)
    if title:
        add_text(slide, title, x + Inches(0.18), cy,
                 w - Inches(0.36), Inches(0.38),
                 size=title_size, bold=True, color=title_color)
        cy += Inches(0.38)
    if body_lines:
        add_multiline(slide, body_lines, x + Inches(0.18), cy,
                      w - Inches(0.36), h - (cy - y) - Inches(0.1),
                      size=body_size, color=body_color)

def flow_box(slide, label, x, y, w=Inches(1.9), h=Inches(0.7), bg=NAVY, fg=WHITE):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, label, x, y, w, h,
             size=10, bold=True, color=fg, align=PP_ALIGN.CENTER)

def arrow(slide, x, y, w=Inches(0.35)):
    add_text(slide, "→", x, y - Inches(0.05), w, Inches(0.5),
             size=18, bold=True, color=GREY, align=PP_ALIGN.CENTER)

def dot_bullet(slide, items, x, y, w, size=12.5, color=DARK, dot_color=NAVY):
    lines = [f"●  {item}" for item in items]
    add_multiline(slide, lines, x, y, w, Inches(len(items) * 0.35 + 0.2),
                  size=size, color=color)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(sl, Inches(0), Inches(0), W, H, fill=NAVY)

# Blue accent diagonal stripe (simulated with a wide rect rotated—use rect instead)
add_rect(sl, Inches(8.5), Inches(0), Inches(4.83), H, fill=BLUE)
add_rect(sl, Inches(8.2), Inches(0), Inches(0.5), H, fill=RGBColor(0x00, 0x8A, 0xB5))

# Paytm logo text
add_text(sl, "Pay", Inches(0.7), Inches(0.5), Inches(3), Inches(1.0),
         size=48, bold=True, color=WHITE)
add_text(sl, "tm", Inches(2.15), Inches(0.5), Inches(3), Inches(1.0),
         size=48, bold=True, color=YELLOW)

# Tagline chip
pill(sl, "  HACKATHON 2025  ", Inches(0.7), Inches(1.45),
     w=Inches(2.4), h=Inches(0.35), bg=YELLOW, fg=NAVY, size=10)

# Main title
add_text(sl, "AI-Powered", Inches(0.7), Inches(2.0), Inches(7.5), Inches(0.9),
         size=52, bold=True, color=WHITE)
add_text(sl, "Offline Payments", Inches(0.7), Inches(2.8), Inches(7.5), Inches(0.9),
         size=52, bold=True, color=YELLOW)

# Subtitle
add_text(sl, "Pay without internet.  Always.", Inches(0.7), Inches(3.75),
         Inches(7.5), Inches(0.6), size=22, color=LIGHT, italic=True)

# Three feature pills
for i, (label, col) in enumerate([
    ("  AI Credit Limit  ", BLUE),
    ("  Bluetooth P2P  ", GREEN),
    ("  Auto-Sync  ", ORANGE),
]):
    pill(sl, label, Inches(0.7 + i * 2.55), Inches(4.5),
         w=Inches(2.3), h=Inches(0.38), bg=col, fg=WHITE, size=10)

# Team
add_text(sl, "Team Issavibles", Inches(0.7), Inches(5.5),
         Inches(5), Inches(0.4), size=14, color=LIGHT)

# Right side graphic placeholders
add_text(sl, "📱", Inches(9.0), Inches(1.5), Inches(2), Inches(2),
         size=96, align=PP_ALIGN.CENTER, color=WHITE)
add_text(sl, "⚡  Offline-First UPI", Inches(8.6), Inches(3.8), Inches(4.0), Inches(0.5),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "The Problem", "India's digital payment system has a critical blind spot")

# Four stat cards
stats = [
    ("66%", "Indians face payment disruptions every month due to poor connectivity"),
    ("500M+", "Users in Tier 2/3 cities where internet is unreliable"),
    ("₹0", "Offline payment capability in any major UPI app today"),
    ("15–30s", "Average UPI timeout — leading to double-payments & confusion"),
]
for i, (num, desc) in enumerate(stats):
    cx = Inches(0.3 + i * 3.26)
    add_rect(sl, cx, Inches(1.4), Inches(3.1), Inches(2.2), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, cx, Inches(1.4), Inches(3.1), Inches(0.06), fill=NAVY)
    add_text(sl, num, cx, Inches(1.5), Inches(3.1), Inches(0.9),
             size=38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_multiline(sl, [desc], cx + Inches(0.15), Inches(2.42),
                  Inches(2.8), Inches(0.9), size=11, color=GREY,
                  align=PP_ALIGN.CENTER)

# Pain points
add_text(sl, "Root Causes", Inches(0.3), Inches(3.85), Inches(12), Inches(0.4),
         size=15, bold=True, color=NAVY)
pain = [
    ("No connectivity = no payment", "Every UPI transaction requires a live bank server call. No fallback exists."),
    ("Merchants bear the cost", "Failed payments at the counter = lost sales, awkward moments, and declining trust."),
    ("Rural & metro dead zones", "Basements, metro stations, mountains, and 2G zones all break the payment flow."),
]
for i, (title, desc) in enumerate(pain):
    cx = Inches(0.3 + i * 4.35)
    add_rect(sl, cx, Inches(4.25), Inches(4.1), Inches(2.7), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, cx, Inches(4.25), Inches(0.06), Inches(2.7), fill=ORANGE)
    add_text(sl, title, cx + Inches(0.15), Inches(4.32),
             Inches(3.8), Inches(0.45), size=13, bold=True, color=NAVY)
    add_text(sl, desc, cx + Inches(0.15), Inches(4.75),
             Inches(3.8), Inches(1.0), size=11, color=GREY)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Our Solution", "Decouple payment capture from settlement")

# Key insight banner
add_rect(sl, Inches(0.3), Inches(1.3), Inches(12.73), Inches(0.7), fill=NAVY)
add_text(sl, "💡  Core Insight: A payment can be captured offline and settled later — safely, with cryptographic guarantees.",
         Inches(0.5), Inches(1.35), Inches(12.5), Inches(0.6),
         size=13, bold=True, color=WHITE)

# Three pillars
pillars = [
    ("1", "AI/ML\nCredit Limit", NAVY,
     ["Server assigns per-user offline limit",
      "Based on: KYC, transaction history,",
      "device trust, fraud signals",
      "Cached on-device for 24 hours",
      "Expires if not refreshed online"]),
    ("2", "Signed Offline\nPayments", BLUE,
     ["Ed25519 cryptographic signatures",
      "Nonce-based replay protection",
      "3 modes: Online, Offline, BLE P2P",
      "Stored in encrypted SQLite queue",
      "Debit local limit immediately"]),
    ("3", "Auto-Sync\nReconciliation", GREEN,
     ["Triggers on connectivity restore",
      "Idempotent batch submission",
      "Fraud check + signature verify",
      "Ledger settlement + new limit",
      "Handles duplicates gracefully"]),
]
for i, (num, title, col, body) in enumerate(pillars):
    cx = Inches(0.3 + i * 4.35)
    add_rect(sl, cx, Inches(2.2), Inches(4.1), Inches(4.8), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, cx, Inches(2.2), Inches(4.1), Inches(0.08), fill=col)
    # Number badge
    add_rect(sl, cx + Inches(0.2), Inches(2.35), Inches(0.45), Inches(0.45), fill=col)
    add_text(sl, num, cx + Inches(0.2), Inches(2.35), Inches(0.45), Inches(0.45),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, cx + Inches(0.75), Inches(2.35),
             Inches(3.2), Inches(0.6), size=14, bold=True, color=col)
    for j, line in enumerate(body):
        add_text(sl, f"•  {line}", cx + Inches(0.2), Inches(3.1 + j * 0.35),
                 Inches(3.7), Inches(0.38), size=11, color=DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "System Architecture", "What runs on-device vs. on-server")

# On-device column
add_rect(sl, Inches(0.3), Inches(1.35), Inches(3.9), Inches(5.7), fill=WHITE,
         line=LIGHT, line_w=Pt(1))
add_rect(sl, Inches(0.3), Inches(1.35), Inches(3.9), Inches(0.45), fill=NAVY)
add_text(sl, "📱  On-Device (Flutter)", Inches(0.35), Inches(1.38), Inches(3.8), Inches(0.4),
         size=12, bold=True, color=WHITE)

device_components = [
    ("Offline Limit Service", "Caches AI limit · expires 24h"),
    ("Payment Blob Queue", "Encrypted SQLite · signed blobs"),
    ("BLE Service", "Advertise + scan + transfer"),
    ("Connectivity Monitor", "Live stream · triggers sync"),
    ("Sync Engine", "Background · idempotent POST"),
    ("Local Risk Scorer", "Penalty per pending payment"),
]
for i, (name, desc) in enumerate(device_components):
    add_rect(sl, Inches(0.45), Inches(1.95 + i * 0.75), Inches(3.6), Inches(0.65),
             fill=LIGHT)
    add_text(sl, name, Inches(0.6), Inches(1.98 + i * 0.75), Inches(3.3), Inches(0.3),
             size=11, bold=True, color=NAVY)
    add_text(sl, desc, Inches(0.6), Inches(2.25 + i * 0.75), Inches(3.3), Inches(0.28),
             size=9.5, color=GREY)

# Arrow
add_text(sl, "⇄", Inches(4.25), Inches(3.9), Inches(0.7), Inches(0.7),
         size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "REST\nAPI", Inches(4.2), Inches(4.55), Inches(0.8), Inches(0.55),
         size=9, color=GREY, align=PP_ALIGN.CENTER)
add_text(sl, "⇄", Inches(4.25), Inches(3.0), Inches(0.7), Inches(0.7),
         size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "BLE", Inches(4.25), Inches(3.65), Inches(0.7), Inches(0.3),
         size=9, color=BLUE, align=PP_ALIGN.CENTER)

# Backend column
add_rect(sl, Inches(5.1), Inches(1.35), Inches(4.0), Inches(5.7), fill=WHITE,
         line=LIGHT, line_w=Pt(1))
add_rect(sl, Inches(5.1), Inches(1.35), Inches(4.0), Inches(0.45), fill=BLUE)
add_text(sl, "☁️  Backend (FastAPI + Render)", Inches(5.15), Inches(1.38), Inches(3.9), Inches(0.4),
         size=12, bold=True, color=WHITE)

backend_components = [
    ("Auth Service", "JWT · secure token management"),
    ("ML Risk Engine", "sklearn · heuristic fallback"),
    ("Token Vault", "Ed25519 signing · expiry"),
    ("Offline Sync API", "Idempotent · batch reconcile"),
    ("Fraud Engine", "Velocity · anomaly · blacklist"),
    ("Ledger Service", "Double-entry · audit trail"),
]
for i, (name, desc) in enumerate(backend_components):
    add_rect(sl, Inches(5.25), Inches(1.95 + i * 0.75), Inches(3.7), Inches(0.65),
             fill=LIGHT)
    add_text(sl, name, Inches(5.4), Inches(1.98 + i * 0.75), Inches(3.4), Inches(0.3),
             size=11, bold=True, color=BLUE)
    add_text(sl, desc, Inches(5.4), Inches(2.25 + i * 0.75), Inches(3.4), Inches(0.28),
             size=9.5, color=GREY)

# DB column
add_rect(sl, Inches(9.2), Inches(1.35), Inches(3.9), Inches(5.7), fill=WHITE,
         line=LIGHT, line_w=Pt(1))
add_rect(sl, Inches(9.2), Inches(1.35), Inches(3.9), Inches(0.45), fill=GREEN)
add_text(sl, "🗄️  Data (PostgreSQL)", Inches(9.25), Inches(1.38), Inches(3.8), Inches(0.4),
         size=12, bold=True, color=WHITE)

db_tables = [
    ("users", "id, balance, kyc_tier, device_trust_score, fraud_flags, offline_limit"),
    ("transactions", "id, sender_id, receiver_id, amount, nonce, status, settled_at"),
    ("offline_tokens", "token_id, user_id, amount, signature, expiry, is_consumed"),
    ("ledger_entries", "user_id, tx_id, entry_type, amount, balance_after"),
    ("payment_blobs", "id, sender_id, receiver_id, nonce, is_offline, status"),
]
for i, (table, cols) in enumerate(db_tables):
    add_rect(sl, Inches(9.35), Inches(1.95 + i * 0.88), Inches(3.6), Inches(0.8), fill=LIGHT)
    add_text(sl, table, Inches(9.5), Inches(1.98 + i * 0.88), Inches(3.3), Inches(0.3),
             size=11, bold=True, color=GREEN)
    add_text(sl, cols, Inches(9.5), Inches(2.25 + i * 0.88), Inches(3.3), Inches(0.45),
             size=8.5, color=GREY)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — THREE PAYMENT CASES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Three Payment Cases", "Covering every connectivity scenario")

cases = [
    {
        "title": "Case 1 — Sender Offline",
        "color": ORANGE,
        "steps": ["Check local\noffline limit", "Create signed\npayment blob", "Store in\nSQLite queue", "Show 'Sent\n(offline)'", "Auto-sync\nwhen online"],
    },
    {
        "title": "Case 2 — Receiver Offline",
        "color": BLUE,
        "steps": ["Sender online\nAPI call", "Debit bank\nimmediately", "Server holds\ncredit", "Receiver\nsyncs", "Credit\napplied"],
    },
    {
        "title": "Case 3 — Both Offline",
        "color": NAVY,
        "steps": ["Scan QR\n(BLE UUID)", "BLE device\ndiscovery", "Blob sent\nvia BLE", "Both queue\nlocally", "First online\nsettles"],
    },
]
for row, case in enumerate(cases):
    cy = Inches(1.45 + row * 1.95)
    col = case["color"]
    # Label
    add_rect(sl, Inches(0.3), cy, Inches(2.1), Inches(1.55), fill=col)
    add_text(sl, case["title"], Inches(0.35), cy + Inches(0.5),
             Inches(2.0), Inches(0.7), size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # Flow boxes
    for i, step in enumerate(case["steps"]):
        bx = Inches(2.6 + i * 2.12)
        add_rect(sl, bx, cy + Inches(0.42), Inches(1.85), Inches(0.72), fill=WHITE,
                 line=col, line_w=Pt(1.5))
        add_text(sl, step, bx, cy + Inches(0.42), Inches(1.85), Inches(0.72),
                 size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text(sl, "→", Inches(2.6 + i * 2.12 + 1.87), cy + Inches(0.56),
                     Inches(0.25), Inches(0.4), size=14, bold=True,
                     color=col, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — AI / ML RISK ENGINE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "AI / ML Risk Engine", "The intelligence behind every offline limit")

# Input features
add_text(sl, "INPUT FEATURES", Inches(0.3), Inches(1.35), Inches(4.0), Inches(0.35),
         size=10, bold=True, color=GREY)
features = [
    ("📊  Transaction Count", "How many payments last 30 days"),
    ("💰  Avg. Transaction Value", "Typical spend pattern"),
    ("🪪  KYC Tier (1–3)", "Verification level"),
    ("📱  Device Trust Score", "Device reputation 0–1"),
    ("📅  Account Age", "Days since registration"),
    ("🚨  Fraud Flags", "Past suspicious activity"),
]
for i, (name, desc) in enumerate(features):
    add_rect(sl, Inches(0.3), Inches(1.75 + i * 0.72), Inches(3.9), Inches(0.65),
             fill=WHITE, line=LIGHT, line_w=Pt(0.5))
    add_text(sl, name, Inches(0.45), Inches(1.78 + i * 0.72), Inches(3.6), Inches(0.3),
             size=11, bold=True, color=NAVY)
    add_text(sl, desc, Inches(0.45), Inches(2.05 + i * 0.72), Inches(3.6), Inches(0.28),
             size=9.5, color=GREY)

# Arrow
add_text(sl, "➜", Inches(4.3), Inches(3.8), Inches(0.6), Inches(0.6),
         size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "sklearn\nmodel", Inches(4.2), Inches(4.4), Inches(0.8), Inches(0.5),
         size=8.5, color=GREY, align=PP_ALIGN.CENTER)

# Risk tiers
add_text(sl, "OFFLINE LIMIT TIERS", Inches(5.0), Inches(1.35), Inches(4.0), Inches(0.35),
         size=10, bold=True, color=GREY)
tiers = [
    ("< 0.2  (Very Low Risk)", "₹5,000", GREEN),
    ("0.2–0.4  (Low Risk)", "₹3,000", GREEN),
    ("0.4–0.6  (Medium Risk)", "₹1,500", ORANGE),
    ("0.6–0.8  (High Risk)", "₹500", ORANGE),
    ("0.8–0.9  (Very High)", "₹100", RGBColor(0xE5, 0x39, 0x35)),
    ("> 0.9  (Blocked)", "₹0", RGBColor(0xE5, 0x39, 0x35)),
]
for i, (label, limit, col) in enumerate(tiers):
    add_rect(sl, Inches(5.0), Inches(1.75 + i * 0.72), Inches(4.3), Inches(0.65),
             fill=WHITE, line=LIGHT, line_w=Pt(0.5))
    add_text(sl, label, Inches(5.15), Inches(1.82 + i * 0.72), Inches(2.9), Inches(0.4),
             size=11, color=DARK)
    add_text(sl, limit, Inches(8.0), Inches(1.82 + i * 0.72), Inches(1.2), Inches(0.4),
             size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)

# Local re-scoring explainer
add_rect(sl, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.7), fill=WHITE,
         line=NAVY, line_w=Pt(1.5))
add_rect(sl, Inches(9.5), Inches(1.35), Inches(3.6), Inches(0.45), fill=NAVY)
add_text(sl, "Local Re-Scoring", Inches(9.6), Inches(1.38), Inches(3.4), Inches(0.4),
         size=12, bold=True, color=WHITE)
lines = [
    "After EACH offline payment:",
    "",
    "penalty = 1 – (pending × 0.1)",
    "capped at 30% floor",
    "",
    "Example:",
    "1 pending → 90% of limit",
    "2 pending → 80% of limit",
    "5 pending → 50% of limit",
    "7+ pending → 30% (floor)",
    "",
    "→ Prevents chaining exploit",
    "→ Runs entirely on-device",
    "→ No server needed",
]
for i, line in enumerate(lines):
    is_code = line.startswith("penalty") or line.startswith("capped")
    col = BLUE if is_code else (NAVY if line.startswith("→") else DARK)
    add_text(sl, line, Inches(9.65), Inches(1.9 + i * 0.3), Inches(3.3), Inches(0.32),
             size=10, bold=is_code or line.startswith("→"), color=col)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — SECURITY ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Security Architecture", "Trust without connectivity")

sec_cards = [
    ("🔐  Ed25519 Signatures", NAVY, [
        "Server holds private key (env var on Render)",
        "Public key distributed to all clients",
        "Every token cryptographically signed",
        "Device verifies before accepting payment",
        "Tamper-evident — any modification = invalid",
    ]),
    ("🔁  Replay Protection", BLUE, [
        "UUID nonce generated per payment blob",
        "Backend deduplicates by (sender+receiver+nonce)",
        "Idempotency key on all sync submissions",
        "Prevents double-spend even if synced twice",
        "SQLite queue tracks blob state locally",
    ]),
    ("🚨  Fraud Detection", ORANGE, [
        "Velocity checks: >5 txns/min = flagged",
        "Amount anomaly: >3× avg = reviewed",
        "Consecutive offline cap (local re-scoring)",
        "Blacklist lookup on every sync",
        "Fraud flags persist across sessions",
    ]),
    ("⏱️  Expiry & Caps", GREEN, [
        "Offline limit expires in 24 hours",
        "Expired + offline = limit becomes ₹0",
        "Max 10 tokens cached per session",
        "Per-session transaction count cap",
        "Dynamic revocation via flag propagation",
    ]),
]
for i, (title, col, body) in enumerate(sec_cards):
    row, col_idx = divmod(i, 2)
    cx = Inches(0.3 + col_idx * 6.5)
    cy = Inches(1.35 + row * 2.9)
    add_rect(sl, cx, cy, Inches(6.2), Inches(2.65), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, cx, cy, Inches(6.2), Inches(0.08), fill=col)
    add_text(sl, title, cx + Inches(0.18), cy + Inches(0.15),
             Inches(5.8), Inches(0.4), size=13, bold=True, color=col)
    for j, line in enumerate(body):
        add_text(sl, f"•  {line}", cx + Inches(0.18), cy + Inches(0.65 + j * 0.37),
                 Inches(5.8), Inches(0.38), size=11, color=DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — APP DEMO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Live App Demo", "Fully functional on physical Android devices — tested offline")

screens = [
    ("Home Screen", "Offline limit badge prominently displayed.\nPromo banner: 'Pay without internet. Always.'", NAVY),
    ("Offline Mode Banner", "Instant orange banner when internet drops.\nAnimated — appears within 1 second.", ORANGE),
    ("Scan & Pay", "Scan any user's QR (merchant, retailer, user).\nAmount input with limit check inline.", BLUE),
    ("Payment Receipt", "Pending Sync status with numbered steps.\nCopyable payment ID for audit trail.", GREEN),
    ("AI Limit Explainer", "Tappable limit shows per-factor breakdown.\nCredit score 0–100 with progress bar.", RGBColor(0x7C, 0x4D, 0xFF)),
    ("My QR Code", "Universal QR — works for all user types.\nBLE UUID embedded for Case 3.", NAVY),
]
for i, (title, desc, col) in enumerate(screens):
    row, c = divmod(i, 3)
    cx = Inches(0.3 + c * 4.35)
    cy = Inches(1.35 + row * 3.0)
    add_rect(sl, cx, cy, Inches(4.1), Inches(2.75), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    # Phone mockup area
    add_rect(sl, cx + Inches(0.15), cy + Inches(0.1), Inches(1.4), Inches(2.2), fill=LIGHT)
    add_rect(sl, cx + Inches(0.15), cy + Inches(0.1), Inches(1.4), Inches(0.4), fill=col)
    add_text(sl, "📱", cx + Inches(0.2), cy + Inches(0.5), Inches(1.3), Inches(1.5),
             size=36, align=PP_ALIGN.CENTER, color=col)
    add_text(sl, title, cx + Inches(1.65), cy + Inches(0.12),
             Inches(2.3), Inches(0.4), size=12, bold=True, color=col)
    add_text(sl, desc, cx + Inches(1.65), cy + Inches(0.55),
             Inches(2.3), Inches(1.4), size=10, color=GREY)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — SYNC & RECONCILIATION ENGINE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Sync & Reconciliation Engine", "What happens when you come back online")

# Workflow nodes (n8n-style horizontal flow)
nodes = [
    ("Internet\nRestored", NAVY),
    ("Background\nSync Trigger", BLUE),
    ("Collect\nPending Blobs", NAVY),
    ("POST\n/api/offline/sync", BLUE),
    ("Validate\nSignatures", ORANGE),
    ("Fraud\nCheck", ORANGE),
    ("Settle\nLedger", GREEN),
    ("Push New\nLimit", NAVY),
]
box_w = Inches(1.4)
box_h = Inches(0.8)
gap   = Inches(0.12)
start_x = Inches(0.3)
cy = Inches(2.1)

for i, (label, col) in enumerate(nodes):
    bx = start_x + i * (box_w + gap)
    add_rect(sl, bx, cy, box_w, box_h, fill=col)
    add_text(sl, label, bx, cy, box_w, box_h,
             size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        ax = bx + box_w
        add_text(sl, "→", ax, cy + Inches(0.15), gap, Inches(0.5),
                 size=11, bold=True, color=GREY, align=PP_ALIGN.CENTER)

# Outcomes
add_text(sl, "RECONCILIATION OUTCOMES", Inches(0.3), Inches(3.25), Inches(12), Inches(0.35),
         size=10, bold=True, color=GREY)
outcomes = [
    ("✅  ACCEPTED", GREEN, "Signature valid · nonce unique · fraud-free · amount within limit\n→ Deduct from user's bank balance · Credit merchant · Update ledger"),
    ("⚠️  ADJUSTED", ORANGE, "Amount exceeds remaining limit (race condition between devices)\n→ Partially settle up to available limit · Return delta to sender"),
    ("❌  REJECTED", RGBColor(0xE5, 0x39, 0x35), "Duplicate nonce · invalid signature · fraud flagged · limit exceeded\n→ Reverse local deduction · Notify user · Flag for review"),
    ("🔁  DUPLICATE", BLUE, "Same (sender+receiver+nonce) already processed by another sync\n→ Mark as settled · No double-credit · Safe to retry"),
]
for i, (title, col, desc) in enumerate(outcomes):
    cx = Inches(0.3 + i * 3.27)
    add_rect(sl, cx, Inches(3.65), Inches(3.1), Inches(3.3), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, cx, Inches(3.65), Inches(3.1), Inches(0.08), fill=col)
    add_text(sl, title, cx + Inches(0.15), Inches(3.78),
             Inches(2.8), Inches(0.38), size=12, bold=True, color=col)
    add_text(sl, desc, cx + Inches(0.15), Inches(4.2),
             Inches(2.8), Inches(2.4), size=10, color=DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — PAYTM INTEGRATION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Integration with Paytm", "Drop-in, not rip-out — additive integration")

# Principle banner
add_rect(sl, Inches(0.3), Inches(1.35), Inches(12.73), Inches(0.65), fill=NAVY)
add_text(sl, "Zero changes to existing payment flow.  Three Flutter services + one backend module behind a feature flag.",
         Inches(0.5), Inches(1.42), Inches(12.5), Inches(0.5),
         size=12, bold=True, color=WHITE)

# Flutter additions
add_text(sl, "FLUTTER (3 services to add)", Inches(0.3), Inches(2.2), Inches(6.0), Inches(0.35),
         size=10, bold=True, color=GREY)
flutter_items = [
    ("OfflineLimitService", "SharedPreferences · limit caching · expiry · local re-scoring"),
    ("OfflineQueueService", "SQLite · blob CRUD · pending/synced state tracking"),
    ("SyncEngine", "Background service · connectivity listener · idempotent POST"),
]
for i, (name, desc) in enumerate(flutter_items):
    add_rect(sl, Inches(0.3), Inches(2.6 + i * 0.85), Inches(5.9), Inches(0.78),
             fill=WHITE, line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, Inches(0.3), Inches(2.6 + i * 0.85), Inches(0.07), Inches(0.78),
             fill=NAVY)
    add_text(sl, name, Inches(0.5), Inches(2.65 + i * 0.85),
             Inches(5.5), Inches(0.32), size=12, bold=True, color=NAVY)
    add_text(sl, desc, Inches(0.5), Inches(2.95 + i * 0.85),
             Inches(5.5), Inches(0.28), size=10, color=GREY)

# Backend additions
add_text(sl, "BACKEND (new routes, same DB)", Inches(6.5), Inches(2.2), Inches(6.5), Inches(0.35),
         size=10, bold=True, color=GREY)
backend_items = [
    ("GET /api/user/offline-limit", "Returns AI limit + risk score + expiry"),
    ("POST /api/offline/sync", "Batch blob reconciliation · per-blob status"),
    ("POST /api/payments/online", "Online payment · handles offline receiver"),
]
for i, (endpoint, desc) in enumerate(backend_items):
    add_rect(sl, Inches(6.5), Inches(2.6 + i * 0.85), Inches(6.5), Inches(0.78),
             fill=WHITE, line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, Inches(6.5), Inches(2.6 + i * 0.85), Inches(0.07), Inches(0.78),
             fill=BLUE)
    add_text(sl, endpoint, Inches(6.7), Inches(2.65 + i * 0.85),
             Inches(6.2), Inches(0.32), size=11, bold=True, color=BLUE,
             font="Courier New")
    add_text(sl, desc, Inches(6.7), Inches(2.95 + i * 0.85),
             Inches(6.2), Inches(0.28), size=10, color=GREY)

# Rollout plan
add_text(sl, "ROLLOUT PLAN", Inches(0.3), Inches(5.1), Inches(12), Inches(0.35),
         size=10, bold=True, color=GREY)
rollout = [
    ("Feature Flag", "Enable for 0% → 1% → 10% → 100% of users"),
    ("Pilot Markets", "Tier 2 cities first — Jaipur, Lucknow, Indore"),
    ("Monitor", "Sync success rate · fraud rate · limit utilisation"),
    ("Iterate", "Retrain ML model monthly with real usage data"),
]
for i, (step, desc) in enumerate(rollout):
    cx = Inches(0.3 + i * 3.27)
    add_rect(sl, cx, Inches(5.5), Inches(3.1), Inches(1.65), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_text(sl, step, cx + Inches(0.15), Inches(5.58),
             Inches(2.8), Inches(0.38), size=12, bold=True, color=NAVY)
    add_text(sl, desc, cx + Inches(0.15), Inches(5.95),
             Inches(2.8), Inches(0.65), size=10, color=GREY)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS IMPACT
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=LIGHT)
section_header(sl, "Business Impact", "Why this matters at Paytm scale")

impact = [
    ("500M+", "Addressable users in Tier 2/3 & rural India — currently underserved by digital payments"),
    ("66%", "Of Indian consumers face monthly payment disruptions — this feature directly resolves each one"),
    ("₹0 → ₹5,000", "Per-user offline limit creates a new class of micro-credit without touching the banking system"),
    ("2 weeks", "Integration timeline — drop-in services, feature flag rollout, no existing code modified"),
]
for i, (num, desc) in enumerate(impact):
    cx = Inches(0.3 + (i % 2) * 6.5)
    cy = Inches(1.35 + (i // 2) * 1.95)
    add_rect(sl, cx, cy, Inches(6.2), Inches(1.75), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_rect(sl, cx, cy, Inches(0.07), Inches(1.75), fill=NAVY)
    add_text(sl, num, cx + Inches(0.25), cy + Inches(0.2),
             Inches(2.5), Inches(0.8), size=44, bold=True, color=NAVY)
    add_text(sl, desc, cx + Inches(0.25), cy + Inches(0.95),
             Inches(5.8), Inches(0.65), size=11, color=GREY)

# Competitive advantage
add_text(sl, "COMPETITIVE MOAT", Inches(0.3), Inches(5.35), Inches(12), Inches(0.35),
         size=10, bold=True, color=GREY)
moat = [
    ("vs. Google Pay", "No offline payments · requires live internet · no BLE fallback"),
    ("vs. PhonePe", "No AI-assigned limit · no P2P offline transfer capability"),
    ("vs. Cash", "Fully digital · auditable · reversible · no change-making required"),
    ("vs. Current UPI", "Additive, not disruptive · compatible with existing NPCI rails"),
]
for i, (vs, desc) in enumerate(moat):
    cx = Inches(0.3 + i * 3.27)
    add_rect(sl, cx, Inches(5.75), Inches(3.1), Inches(1.4), fill=WHITE,
             line=LIGHT, line_w=Pt(0.5))
    add_text(sl, vs, cx + Inches(0.15), Inches(5.83),
             Inches(2.8), Inches(0.38), size=11, bold=True, color=NAVY)
    add_text(sl, desc, cx + Inches(0.15), Inches(6.18),
             Inches(2.8), Inches(0.72), size=9.5, color=GREY)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — ROADMAP & TEAM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, Inches(0), Inches(0), W, H, fill=NAVY)

# Title
add_text(sl, "Roadmap & What's Next", Inches(0.5), Inches(0.3),
         Inches(10), Inches(0.7), size=36, bold=True, color=WHITE)
add_text(sl, "A hackathon prototype → a production-ready feature",
         Inches(0.5), Inches(0.95), Inches(10), Inches(0.45),
         size=14, color=LIGHT, italic=True)

# Timeline phases
phases = [
    ("Phase 1\n0–3 Months", BLUE, [
        "Pilot: 1,000 merchants in Jaipur",
        "A/B test offline limit tiers",
        "Collect sync failure telemetry",
        "Refine ML model with real data",
        "Security audit",
    ]),
    ("Phase 2\n3–6 Months", YELLOW, [
        "Full UPI stack integration",
        "NPCI offline payment sandbox",
        "iOS + Android GA release",
        "Merchant dashboard analytics",
        "Fraud model v2 training",
    ]),
    ("Phase 3\n6–12 Months", GREEN, [
        "National rollout (all Paytm users)",
        "RBI regulatory compliance",
        "Cross-bank offline payments",
        "International expansion",
        "NFC + Bluetooth 5.0 upgrade",
    ]),
]
for i, (title, col, items) in enumerate(phases):
    cx = Inches(0.4 + i * 4.3)
    add_rect(sl, cx, Inches(1.6), Inches(4.0), Inches(4.5), fill=RGBColor(0x0A, 0x1A, 0x4A))
    add_rect(sl, cx, Inches(1.6), Inches(4.0), Inches(0.75), fill=col)
    add_text(sl, title, cx, Inches(1.6), Inches(4.0), Inches(0.75),
             size=14, bold=True, color=NAVY if col == YELLOW else WHITE,
             align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, f"•  {item}", cx + Inches(0.2), Inches(2.5 + j * 0.52),
                 Inches(3.6), Inches(0.45), size=11, color=LIGHT)

# Team & ask
add_rect(sl, Inches(0.4), Inches(6.15), Inches(12.53), Inches(1.05),
         fill=RGBColor(0x0A, 0x1A, 0x4A))
add_text(sl, "🏆  Team Issavibles", Inches(0.6), Inches(6.22),
         Inches(4), Inches(0.4), size=14, bold=True, color=YELLOW)
add_text(sl, "Built & shipped in <48 hours: Flutter (iOS + Android) · FastAPI · PostgreSQL · BLE · ML · Deployed on Render",
         Inches(0.6), Inches(6.58), Inches(7), Inches(0.45),
         size=11, color=LIGHT)
add_text(sl, "Ask: Paytm merchant network for pilot · Backend infra support",
         Inches(7.8), Inches(6.22), Inches(5), Inches(0.4),
         size=12, bold=True, color=YELLOW)
add_text(sl, "paywithoutwifi@issavibles.com",
         Inches(7.8), Inches(6.6), Inches(5), Inches(0.35),
         size=11, color=LIGHT)

# ── Save ───────────────────────────────────────────────────────
out = "/Users/vivekgupta/payapp_/PaytmOfflinePay_Hackathon.pptx"
prs.save(out)
print(f"Saved: {out}")
